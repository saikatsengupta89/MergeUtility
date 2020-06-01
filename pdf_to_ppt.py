import os
import subprocess
import fitz as fz
import PyPDF2
from pptx import Presentation
from PIL import Image
from io import BytesIO

from xlrd.formatting import fmt_bracketed_sub
pdf_dir=r"C:\MergePDF\Output"
os.chdir(pdf_dir)
lst=[os.listdir(pdf_dir)]
print(lst)
pdffile = "ConsolidatedGroupPack_c.pdf"
pdfFile = open(pdffile, 'rb')
pdfReader = PyPDF2.PdfFileReader(pdfFile)
print(pdfReader.getNumPages())
lst_images= []
zoom=50/72 #controls image resolution
mat= fz.Matrix(zoom, zoom)
for p in range(pdfReader.getNumPages()-1):
	doc = fz.open(pdffile)
	page = doc.loadPage(p)  # number of page
	pix = page.getPixmap(matrix=mat, alpha=False)
	output = "outfile_"+str(p)+".png"
	lst_images.append(output)
	pix.writePNG("C:/MergePDF/Images/"+output)
	pix=None

os.chdir("C:/MergePDF/Images")
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
for output in lst_images:
	image = Image.open(output)
	width, height = image.size
	prs.slide_height = height * 9525
	prs.slide_width = width * 9525
	slide = prs.slides.add_slide(blank_slide_layout)
	pic = slide.shapes.add_picture(output, 0, 0, width=width * 9525, height=height * 9525)

prs.save('C:/MergePDF/Output/CompletePack.pptx')


#check compression
open_doc= fz.open("C:/MergePDF/Output/ConsolidatedGroupPack_c.pdf")
open_doc.save("C:/MergePDF/Output/compressed.pdf", garbage=4, deflate=1, linear=1)