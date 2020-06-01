"""
TO BUILD THE EXE USING THE EXISTING SPEC FILE
C:\PycharmProjects\WorkWithPDF\venv\Scripts\pyinstaller.exe --onefile --icon=favicon.ico MergePDFS_v0.3.spec
"""
from tkinter import Message
from tkinter import messagebox
from tkinter import filedialog
from string import ascii_uppercase
from xlrd import open_workbook
from os import path, listdir, chdir, remove, mkdir
import tkinter as tk
import PyPDF2
import re
import shutil
import io
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import locale
import fitz as fz
from pptx import Presentation
from PIL import Image

class Root(tk.Tk):
    dataDrivenFilePath=''
    convPdfFilePath=''

    def __init__ (self):
        super(Root, self).__init__()
        self.title("Merge PDF Utility")
        self.minsize(750, 600)
        self.maxsize(750, 600)
        self.wm_iconbitmap("favicon.ico")


class CustomMessage (tk.Tk):
    def __init__ (self):
        super(CustomMessage, self).__init__()
        self.title("Information")
        self.minsize(500, 100)
        self.maxsize(500, 100)

def atoi(text):
    return int(text) if text.isdigit() else text

def natural_keys(text):
    return [atoi(c) for c in re.split(r'(\d+)', text)]

def convertPdfToPpt(pptFileName):
    #print("Convert PDF Path: " + master.convPdfFilePath)
    #print(path.dirname(master.convPdfFilePath))
    fileName= pptFileName.get()
    dirName=  path.dirname(master.convPdfFilePath)
    mkdir(dirName + "/Images/")

    pdfFile = open(master.convPdfFilePath, 'rb')
    pdfReader = PyPDF2.PdfFileReader(pdfFile)
    #print(pdfReader.getNumPages())

    lst_images = []
    zoom = 100 / 72  # controls image resolution
    mat = fz.Matrix(zoom, zoom)
    for p in range(pdfReader.getNumPages()):
        doc = fz.open(master.convPdfFilePath)
        page = doc.loadPage(p)  # number of page
        pix = page.getPixmap(matrix=mat, alpha=False)
        output = "outfile_" + str(p) + ".png"
        lst_images.append(output)
        pix.writePNG(dirName+"/Images/" + output)
        pix = None

    chdir(dirName+"/Images")
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for output in lst_images:
        image = Image.open(output)
        width, height = image.size
        prs.slide_height = height * 9525
        prs.slide_width = width * 9525
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(output, 0, 0, width=width * 9525, height=height * 9525)

    image.close()
    chdir(dirName)
    prs.save(dirName+"/"+fileName+".pptx")
    shutil.rmtree(dirName+"/Images")

    m = Message(master=CustomMessage(),
                width=500,
                pady=30,
                anchor='center',
                font='bold',
                text="PPT File placed at: " + str(dirName+"/"+fileName+".pptx"))
    m.pack()


def compressOutputPDF(outputFile_plain, outputFile_comp):
    outputFilePath = outputFile_comp
    content = outputFile_plain
    # "-dCompatibilityLevel=1.4",
    # "-dPDFSETTINGS=/ebook",
    args = [
        "ps2pdf",  # actual value doesn't matter
        "-dNOPAUSE", "-dBATCH", "-dSAFER",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        "-dPDFSETTINGS=/screen",
        "-dEmbedAllFonts=true",
        "-dSubsetFonts=true",
        "-dColorImageDownsampleType=/Bicubic",
        "-dColorImageResolution=110",
        "-dGrayImageDownsampleType=/Bicubic",
        "-dGrayImageResolution=110",
        "-sOutputFile=" + outputFilePath,
        "-c", ".setpdfwrite",
        "-f", content
    ]
    # arguments have to be bytes, encode them
    encoding = locale.getpreferredencoding()
    args = [a.encode(encoding) for a in args]

    print(args)
    #import ghostscript
    #ghostscript.Ghostscript(*args)

    #remove the uncompressed file from the path
    # if path.exists(outputFile_plain):
    #     remove(outputFile_plain)

def createPageWithNumbers (pageNum):
    packet = io.BytesIO()
    # create a new PDF with Reportlab
    can = canvas.Canvas(packet, pagesize=letter)
    can.setFontSize(8)
    can.drawString(770, 10, str(pageNum))
    can.save()
    return packet

def putPageNumbers(tempFile, outputFile_plain, outputFile_comp):
    # read your existing PDF
    open_temp= open(tempFile, 'rb')
    existing_pdf = PyPDF2.PdfFileReader(open_temp)
    output = PyPDF2.PdfFileWriter()
    # add the "watermark" (which is the new pdf) on the existing page
    for pagenum in range(existing_pdf.numPages):
        pageObj = existing_pdf.getPage(pagenum)
        # pdfWriter.addPage(pageObj)
        page = existing_pdf.getPage(pagenum)
        if (pagenum > 2):
            packet = createPageWithNumbers(pagenum + 1)
        else:
            packet= createPageWithNumbers('')
        # move to the beginning of the StringIO buffer
        packet.seek(0)
        new_pdf = PyPDF2.PdfFileReader(packet)
        page.mergePage(new_pdf.getPage(0))
        output.addPage(page)

    # finally, write "output" to a real file
    outputStream = open(outputFile_plain, "wb")
    output.write(outputStream)
    outputStream.close()
    open_temp.close()

    #remove the temporary file created during the process using the os.remove() method
    if (path.exists(tempFile)):
        remove(tempFile)

    #not calling compression as ghostscript is not installed in user pc
    #compressOutputPDF(outputFile_plain, outputFile_comp)
    doc= fz.open(outputFile_plain)
    doc.save(outputFile_comp, garbage=4, deflate=1, linear=1)
    doc.close()

    #remove the uncompressed file from the path
    if path.exists(outputFile_plain):
        remove(outputFile_plain)


def mergePDF(dataDrivenFileLoc, inputLocation, outputLocation, outputFileName):
    # move to the input directory and create a list of all pdfs needs to be merged
    chdir(inputLocation)
    file_location= dataDrivenFileLoc
    # wb = load_workbook(inputLocation + "/" + 'InputOrder.xlsx')
    # ws = wb.get_sheet_by_name('InputOrder')
    # column_order = ws['A']
    # column = ws['B']
    workbook= open_workbook(file_location)
    worksheet= workbook.sheet_by_index(0)
    report_order = worksheet.col(0)
    report_name = worksheet.col(1)
    report_type= worksheet.col(2)
    list_master = []
    list_present = []
    list_present_org = []
    list_absent = []
    list_pdf = []
    print(report_order[0])
    print(report_name[0])
    print(report_type[0])
    for row in range(1, len(report_order)):
        if (str(report_type[row].value).upper().strip()=='PDF'):
            list_master.append(str(report_name[row].value).strip().lower() + ".pdf")
    #print(list_master)

    for filename in listdir('.'):
        if filename.endswith('.pdf'):
            list_present.append(filename.lower())
            list_present_org.append(filename)
    #print(list_present)

    for data in list_master:
        if data in list_present:
            for org in list_present_org:
                if (org.lower()==data):
                    list_pdf.append(org)
    #print(list_pdf)

    for data in list_master:
        if data not in list_present:
            list_absent.append(data)

    pdfWriter = PyPDF2.PdfFileWriter()
    # loop through all the pdfs and merge them one by one
    for filename in list_pdf:
        # rb for read binary format
        pdfFile = open(filename, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFile)
        # opening each page in the pdf
        for pagenum in range(pdfReader.numPages):
            pageObj = pdfReader.getPage(pagenum)
            pageObj.compressContentStreams()
            pdfWriter.addPage(pageObj)

    # save the Output in a file, wb for write binary
    tempFile= str(outputLocation +"/"+ outputFileName +"_tmp.pdf")
    outputFile_plain = str(outputLocation + "/" + outputFileName + "_c.pdf")
    outputFile_comp = str(outputLocation + "/" + outputFileName + ".pdf")
    pdfOutput = open(tempFile, 'wb')
    pdfWriter.write(pdfOutput)
    # close the pdfWriter post writing
    pdfOutput.close()

    #call function to write pagenumbers in pdf
    putPageNumbers(tempFile, outputFile_plain, outputFile_comp)

    #to print out the list of PDFs which are not present from the given list of files in the excel
    print("Report Names Missing: ", sep="\n")
    for data in list_absent:
        print (data, sep="\n")


def moveMergePDF(sourcePath, destinationPath, filename):
    # move the output to the network destination folder
    dest_path = destinationPath
    source_path = sourcePath.replace('/','\\')
    file_name = "\\"+filename+".pdf"
    shutil.copyfile(source_path + file_name, dest_path + file_name)

def show_value(ent, dataDrivenFileLoc):
    inputLocation = str(ent['inputFolder'].get())
    outputLocation= str(ent['outputFolder'].get())
    outputFileName= str(ent['fileName'].get())

    tempFile= mergePDF(dataDrivenFileLoc, inputLocation, outputLocation, outputFileName)
    #print("Output Merged File generated as :" + str(outputLocation + "/" + outputFileName + ".pdf"))

    m= Message(master=CustomMessage(),
               width=500,
               pady=30,
               anchor='center',
               font='bold',
               text="Merged File placed at: " + str(outputLocation + "/" + outputFileName + ".pdf"))
    m.pack()


def check_empty(ent):
    dataDrivenFileLoc = master.dataDrivenFilePath
    inputLocation     = str(ent['inputFolder'].get())
    outputLocation    = str(ent['outputFolder'].get())
    outputFileName    = str(ent['fileName'].get())
    available_drives  = ['%s:' % d for d in ascii_uppercase if path.exists('%s:' % d)]
    if (len(dataDrivenFileLoc) > 0 and
        len(inputLocation) > 0 and
        len(outputLocation) > 0 and
        len(outputFileName) > 0):

        if dataDrivenFileLoc =='':
            messagebox.showwarning("Warning", "Data Driven File path doesn't exist.")
        elif path.realpath(inputLocation).replace('\\','').upper() in available_drives:
            messagebox.showwarning("Warning", "Input directory does not exist. Enter proper path.")
        elif path.realpath(outputLocation).replace('\\','').upper() in available_drives:
            messagebox.showwarning("Warning", "Output directory does not exist. Enter proper path.")
        elif not path.isdir(inputLocation):
            messagebox.showwarning("Warning", "Input directory does not exist. Enter proper path.")
        elif not path.isdir(outputLocation):
            messagebox.showwarning("Warning", "Output directory does not exist. Enter proper path.")
        else:
            show_value(ent, dataDrivenFileLoc)
    else:
        messagebox.showwarning("Warning","You must enter all the fields to proceed")

def fileDialog():
    fileName= filedialog.askopenfilename(initialdir= "/", title= "Select a File",
                                              filetype=(("xlsx", "*.xlsx"), ("csv","*.csv")))
    inputFile= tk.Label(text="")
    inputFile.place(x=340, y=145)
    inputFile.configure(text= fileName)
    master.dataDrivenFilePath = fileName
    #print("FilePath: "+master.dataDrivenFilePath)

def fileDialog2():
    fileName= filedialog.askopenfilename(initialdir= "/", title= "Select a File",
                                              filetype=(("xlsx", "*.xlsx"), ("csv","*.csv"),("ALL File", "*")))
    inputFile= tk.Label(text="")
    inputFile.place(x=340, y=435)
    inputFile.configure(text= fileName)
    master.convPdfFilePath = fileName
    #print("FilePath: "+master.dataDrivenFilePath)

def make_form(master):
    ent=dict()
    Label_0 = tk.Label(master, text="MERGE PDFS", width=20, font=("bold", 30))
    Label_0.place(x=120, y=53)

    inputFile = tk.Label(master, text="Data Driven File Location", width=20, font=("bold", 10), anchor='w', justify='left')
    inputFile.place(x=80, y=140)
    button1= tk.Button(text='Browse A File', bg= '#9966ff', fg= '#ffffff', command=lambda: fileDialog())
    button1.place(x=240, y=140)
    #print("From Master: "+master.dataDrivenFilePath)

    Label_1 = tk.Label(master, text="Input Folder Location", width=20, font=("bold", 10), anchor='w', justify='left')
    Label_1.place(x=80, y=190)
    inputFolder= tk.Entry(master, width=70)
    inputFolder.place(x=240, y=190)
    instruction1 = tk.Label(master, text="Example: C:/MergePDF/Input", width=30, font=("normal", 8), anchor='w', justify='left')
    instruction1.place(x=240, y=210)
    ent['inputFolder']=inputFolder

    Label_2 = tk.Label(master, text="Output Folder Location", width=20, font=("bold", 10), anchor='w',
                       justify='left')
    Label_2.place(x=80, y=240)
    outputFolder = tk.Entry(master, width=70)
    outputFolder.place(x=240, y=240)
    instruction2 = tk.Label(master, text="Example: C:/MergePDF/Output", width=30, font=("normal", 8), anchor='w',
                            justify='left')
    instruction2.place(x=240, y=260)
    ent['outputFolder'] = outputFolder

    Label_3 = tk.Label(master, text="Merged File Name", width=20, font=("bold", 10), anchor='w', justify='left')
    Label_3.place(x=80, y=290)
    fileName = tk.Entry(master, width=50)
    fileName.place(x=240, y=290)
    instruction3 = tk.Label(master, text="Specify Desired Output File Name", width=30, font=("normal", 8), anchor='w',
                            justify='left')
    instruction3.place(x=240, y=310)
    ent['fileName'] = fileName

    button2 = tk.Button(master, text='Merge PDFS', width=20, bg='brown', fg='white', command=lambda:check_empty(ent))
    button2.place(x=240, y=340)

    header = tk.Label(master, text="CONVERT PDF TO PPT", width=20, font=("bold, underline", 10), anchor='w',
                         justify='left')
    header .place(x=80, y=400)

    inputFile = tk.Label(master, text="Provide PDF File Location", width=20, font=("bold", 10), anchor='w',
                         justify='left')
    inputFile.place(x=80, y=430)
    button3 = tk.Button(text='Browse A File', bg='#9966ff', fg='#ffffff', command=lambda: fileDialog2())
    button3.place(x=240, y=425)

    Label_4 = tk.Label(master, text="PPT File Name", width=20, font=("bold", 10), anchor='w', justify='left')
    Label_4.place(x=80, y=460)
    pptFileName = tk.Entry(master, width=50)
    pptFileName.place(x=240, y=460)
    instruction3 = tk.Label(master, text="Specify Desired Output File Name", width=30, font=("normal", 8), anchor='w',
                            justify='left')
    instruction3.place(x=240, y=480)
    ent['pptFileName'] = pptFileName

    button4 = tk.Button(master, text='Generate PPT', width=20, bg='brown', fg='white', command=lambda: convertPdfToPpt(pptFileName))
    button4.place(x=240, y=510)

    button5 = tk.Button(master, text='QUIT', width=10, bg='brown', fg='white', command=lambda: master.quit())
    button5.place(x=80, y=530)


if __name__=="__main__":
    master = Root()
    make_form(master)
    master.mainloop()